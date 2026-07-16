# -*- coding: utf-8 -*-
"""文件内容提取：PPTX / PDF / DOCX / 图片 -> 统一 Record。

关键原则：
- 提取阶段无上下文限制，必须把全文逐页落地，绝不截断。
- 每条 Record 带来源溯源（文件名 + 页码），供后续真题索引与回溯。
- 可选依赖缺失时优雅降级（跳过该格式并告警），不中断整批处理。
"""

from __future__ import annotations
import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Record:
    source_file: str          # 文件名（不含路径）
    source_type: str          # pptx | pdf | docx | image
    page: int                 # 页码 / 幻灯片号 / 图片序号
    content: str              # 正文
    title: str = ""           # 所属标题（docx 标题链 / ppt 幻灯片标题）
    heading_level: int = 0    # docx 标题层级（0=正文，1/2/3=Heading）；pptx 标题记 0


# ----------------------------- PPTX -----------------------------
def extract_pptx(path: Path):
    try:
        from pptx import Presentation
    except ImportError:
        print(f"[warn] 未安装 python-pptx，跳过 {path.name}")
        return []
    out = []
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"[warn] 读取 PPTX 失败 {path.name}: {e}")
        return []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        bodies = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            # 识别标题占位符（idx==0 或 shapes.title）
            is_title = False
            try:
                if shape == slide.shapes.title:
                    is_title = True
            except Exception:
                pass
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    is_title = True
            except Exception:
                pass
            if is_title and not title:
                title = txt
            else:
                bodies.append(txt)
        # 演讲者备注——常含正文没有的完整解释
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        content = "\n".join(b for b in bodies if b)
        if notes:
            content += ("\n\n[备注] " + notes) if content else ("[备注] " + notes)
        if title or content:
            out.append(Record(
                source_file=path.name, source_type="pptx", page=i,
                content=content or title, title=title,
            ))
    return out


# ----------------------------- PDF -----------------------------
def extract_pdf(path: Path):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print(f"[warn] 未安装 PyMuPDF，跳过 {path.name}")
        return []
    out = []
    try:
        doc = fitz.open(str(path))
    except Exception as e:
        print(f"[warn] 读取 PDF 失败 {path.name}: {e}")
        return []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        # 扫描页兜底：文本极少则尝试整页 OCR
        if len(text) < 15:
            ocr_text = ""
            try:
                pix = page.get_pixmap(dpi=200)
                import io
                img_path = path.parent / f"__tmp_p{i}_{path.stem}.png"
                pix.save(str(img_path))
                ocr_text = _ocr_image(img_path)
                img_path.unlink(missing_ok=True)
            except Exception:
                pass
            text = ocr_text or text
        if text:
            out.append(Record(
                source_file=path.name, source_type="pdf", page=i, content=text,
            ))
    doc.close()
    return out


# ----------------------------- DOCX -----------------------------
def extract_docx(path: Path):
    try:
        from docx import Document
    except ImportError:
        print(f"[warn] 未安装 python-docx，跳过 {path.name}")
        return []
    out = []
    try:
        doc = Document(str(path))
    except Exception as e:
        print(f"[warn] 读取 DOCX 失败 {path.name}: {e}")
        return []
    heading_chain = []  # [(level, text)]
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        m = re.match(r"heading\s*(\d+)", style)
        if m:
            level = int(m.group(1))
            heading_chain = [(lv, t) for lv, t in heading_chain if lv < level]
            heading_chain.append((level, text))
            chain_str = " > ".join(t for _, t in heading_chain)
            out.append(Record(
                source_file=path.name, source_type="docx", page=0,
                content=text, title=chain_str, heading_level=level,
            ))
        else:
            chain_str = " > ".join(t for _, t in heading_chain)
            out.append(Record(
                source_file=path.name, source_type="docx", page=0,
                content=text, title=chain_str, heading_level=0,
            ))
    # 表格
    for ti, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            out.append(Record(
                source_file=path.name, source_type="docx", page=0,
                content="[表格]\n" + "\n".join(rows),
                title=f"表格{ti}",
            ))
    return out


# ----------------------------- 图片 OCR -----------------------------
_OCR_ENGINE = None


def _get_ocr():
    global _OCR_ENGINE
    if _OCR_ENGINE is not None:
        return _OCR_ENGINE
    try:
        from rapidocr_onnxruntime import RapidOCR
        _OCR_ENGINE = RapidOCR()
    except Exception:
        _OCR_ENGINE = False
    return _OCR_ENGINE


def _ocr_image(path: Path) -> str:
    ocr = _get_ocr()
    if not ocr:
        return ""
    try:
        result, _ = ocr(str(path))
    except Exception:
        return ""
    # RapidOCR 返回 [[box, text, score], ...]，逐项取 text
    return "\n".join(item[1] for item in (result or [])).strip()


def extract_image(path: Path, idx: int):
    text = _ocr_image(path)
    if text:
        return [Record(source_file=path.name, source_type="image", page=idx, content=text)]
    print(f"[warn] 图片 OCR 无文本或未启用 RapidOCR：{path.name}")
    return []


# ----------------------------- 入口 -----------------------------
def extract_dir(input_dir: Path, exclude=None):
    """提取目录下全部支持格式的文件，返回 records。

    exclude: 文件名集合，其中的文件将被跳过（GUI「排除」用——不删原文件，仅不参与生成）。
    """
    exclude = set(exclude or ())
    records = []
    img_idx = 0
    ext_map = {
        ".pptx": extract_pptx, ".ppt": extract_pptx,
        ".pdf": extract_pdf,
        ".docx": extract_docx, ".doc": extract_docx,
    }
    img_exts = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}
    for path in sorted(p for p in input_dir.iterdir() if p.is_file()):
        if path.name in exclude:
            continue
        ext = path.suffix.lower()
        if ext in ext_map:
            records.extend(ext_map[ext](path))
        elif ext in img_exts:
            img_idx += 1
            records.extend(extract_image(path, img_idx))
        else:
            print(f"[skip] 不支持的格式：{path.name}")
    return records
