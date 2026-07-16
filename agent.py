# -*- coding: utf-8 -*-
"""复习笔记智能助手 - 入口（CLI + 样例生成）。

核心流水线见 pipeline.py（5 阶段：分类 -> 课件建骨 -> 书/笔记精修 -> 试题校准 -> 补例题）。

用法：
    py agent.py --make-sample-full _demo        # 生成 4 类齐全的样例
    py agent.py _demo -o 笔记.md --no-llm       # Demo 模式跑通 5 阶段
    py agent.py 你的资料目录 -o 笔记.md --course 课程名   # 真实整理（先 py setup.py 配 API）
"""
from __future__ import annotations
import argparse
import os
from pathlib import Path

import config
import pipeline


def main():
    config.load_env()  # 启动时加载 .env（真实环境变量优先）
    ap = argparse.ArgumentParser(description="复习笔记智能助手（多文件分阶段）")
    ap.add_argument("input_dir", nargs="?", help="学习资料所在目录")
    ap.add_argument("-o", "--output", default="output_notes/复习笔记.md",
                    help="输出 Markdown 路径（默认落 output_notes/ 子文件夹）")
    ap.add_argument("--course", default=os.getenv("REVIEW_COURSE", "未命名课程"),
                    help="课程名称（影响学科自动判定；可用 .env REVIEW_COURSE 设默认）")
    ap.add_argument("--discipline", default=os.getenv("REVIEW_DISCIPLINE", "auto"),
                    choices=["auto", "general", "science", "cs", "business", "arts"],
                    help="学科类型，auto 时按课程名推断")
    ap.add_argument("--chapters", default="", help="显式章节列表，逗号分隔，如 '第1章 A,第2章 B'")
    ap.add_argument("--max-chars", type=int, default=6000, help="单分片最大字符数（超则切）")
    ap.add_argument("--api-base", default="", help="OpenAI 兼容接口 base_url")
    ap.add_argument("--api-key", default="", help="API key")
    ap.add_argument("--model", default="", help="模型名")
    ap.add_argument("--no-llm", action="store_true", help="Demo 模式，不调用大模型")
    ap.add_argument("--no-cache", action="store_true",
                    help="禁用 LLM 响应缓存（强制全重跑；换模型/改 prompt 后用）")
    ap.add_argument("--no-pdf", action="store_true",
                    help="不生成 PDF；默认在产出 MD 后自动转 PDF（需 xelatex + tcolorbox）")
    ap.add_argument("--make-sample", metavar="DIR", help="生成 tiny 样例（课件+笔记）到该目录后退出")
    ap.add_argument("--make-sample-full", metavar="DIR",
                    help="生成 4 类齐全样例（课件/书/笔记/试题）到该目录后退出")
    args = ap.parse_args()
    # PDF 默认自动生成，除非 --no-pdf 或 .env 设了 REVIEW_NO_PDF=1
    args.pdf = not (args.no_pdf or os.getenv("REVIEW_NO_PDF") == "1")

    if args.make_sample_full:
        make_sample_full(Path(args.make_sample_full))
        return
    if args.make_sample:
        make_sample(Path(args.make_sample))
        return
    if not args.input_dir:
        ap.error("请提供输入目录，或用 --make-sample / --make-sample-full 生成样例。")

    run(Path(args.input_dir), Path(args.output), args)


def run(input_dir: Path, out_path: Path, args):
    """薄包装，委托给 pipeline。"""
    pipeline.run_pipeline(input_dir, out_path, args)


# ============================ 样例生成 ============================
def make_sample(out_dir: Path):
    """tiny 样例（课件 pptx + 笔记 docx），用于快速体验。"""
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        from pptx import Presentation
        prs = Presentation()
        lay = prs.slide_layouts[1]
        s = prs.slides.add_slide(lay)
        s.shapes.title.text = "第1章 国民收入核算"
        s.placeholders[1].text = "GDP 是一定时期内一国境内生产的所有最终产品和服务的市场价值总和。"
        s2 = prs.slides.add_slide(lay)
        s2.shapes.title.text = "1.1 GDP 的核算方法"
        s2.placeholders[1].text = "支出法：GDP = C + I + G + NX。\n收入法：把各生产要素收入加总。\n易错：二手交易不计入 GDP。"
        s2.notes_slide.notes_text_frame.text = "教师强调：名义 GDP 与实际 GDP 的区别是高频考点。"
        prs.save(str(out_dir / "宏观经济学.pptx"))
        print(f"[sample] 已生成 {out_dir/'宏观经济学.pptx'}")
    except Exception as e:
        print(f"[sample] 跳过 pptx（{e}）")
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("第2章 总需求与总供给", level=1)
        doc.add_paragraph("总需求曲线反映价格水平与总需求量的反向关系。")
        doc.add_heading("2.1 总需求曲线的移动", level=2)
        doc.add_paragraph("财政政策与货币政策会使 AD 曲线移动。这是论述题高频考点。")
        doc.add_heading("第3章 IS-LM 模型", level=1)
        doc.add_paragraph("IS 曲线表示产品市场均衡，LM 曲线表示货币市场均衡。")
        doc.save(str(out_dir / "宏观经济学笔记.docx"))
        print(f"[sample] 已生成 {out_dir/'宏观经济学笔记.docx'}")
    except Exception as e:
        print(f"[sample] 跳过 docx（{e}）")
    print(f"[sample] 完成。可运行：py agent.py {out_dir} -o 笔记.md --course 宏观经济学")


def make_sample_full(out_dir: Path):
    """4 类齐全样例（课件/书/笔记/试题），用于端到端验证 5 阶段流水线。"""
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1) 课件 PPTX
    try:
        from pptx import Presentation
        prs = Presentation()
        lay = prs.slide_layouts[1]
        s = prs.slides.add_slide(lay)
        s.shapes.title.text = "第1章 供给与需求"
        s.placeholders[1].text = (
            "需求定律：价格上升，需求量下降。\n"
            "供给定律：价格上升，供给量上升。\n"
            "均衡价格由需求曲线与供给曲线的交点决定。"
        )
        s2 = prs.slides.add_slide(lay)
        s2.shapes.title.text = "1.1 弹性"
        s2.placeholders[1].text = (
            "需求价格弹性 Ed = (需求量变化率) / (价格变化率)。\n"
            "|Ed|>1 富有弹性；|Ed|<1 缺乏弹性。\n"
            "例题：若价格下降 10% 使需求量上升 20%，求 Ed。"
        )
        s2.notes_slide.notes_text_frame.text = "教师强调：弹性的计算与分类是高频考点。"
        prs.save(str(out_dir / "01_课件_微观经济学.pptx"))
        print(f"[sample-full] 已生成课件 {out_dir/'01_课件_微观经济学.pptx'}")
    except Exception as e:
        print(f"[sample-full] 跳过 pptx（{e}）")

    # 2) 书籍资料 DOCX（大段说明）
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("第1章 供给与需求", level=1)
        doc.add_paragraph(
            "本章系统讨论市场机制的核心。需求指消费者在一定时期内，在各种价格水平下愿意且能够购买的某种商品的数量。"
            "需求曲线向右下方倾斜，这一性质可由替代效应与收入效应共同解释。供给则指生产者在各种价格水平下愿意且能够出售的数量，"
            "供给曲线向右上方倾斜。当市场需求量等于市场供给量时，市场达到均衡，此时的价格为均衡价格。"
        )
        doc.add_paragraph(
            "应当注意，均衡本身会随需求或供给的变动而移动。例如收入增加使需求曲线右移，导致均衡价格与均衡数量同时上升。"
            "这一比较静态分析是理解市场变化的关键。"
        )
        doc.save(str(out_dir / "02_教材_西方经济学.docx"))
        print(f"[sample-full] 已生成书籍 {out_dir/'02_教材_西方经济学.docx'}")
    except Exception as e:
        print(f"[sample-full] 跳过 book docx（{e}）")

    # 3) 笔记 DOCX（精简要点）
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("复习笔记 - 供给与需求", level=1)
        doc.add_paragraph("重点：需求/供给定律、均衡、弹性的计算。不考：经济学说史背景。")
        doc.add_paragraph("易错：需求量变动（沿曲线移动）与需求变动（曲线本身移动）要分清。")
        doc.add_paragraph("记忆：价升量降（需求）；价升量升（供给）。")
        doc.save(str(out_dir / "03_笔记_重点.docx"))
        print(f"[sample-full] 已生成笔记 {out_dir/'03_笔记_重点.docx'}")
    except Exception as e:
        print(f"[sample-full] 跳过 note docx（{e}）")

    # 4) 试题 DOCX
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("微观经济学 期末试题（往年真题）", level=1)
        doc.add_paragraph("一、选择题（每题 5 分）")
        doc.add_paragraph("1. 某商品价格下降 10%，需求量上升 5%，则其需求价格弹性为：A.0.5 B.2 C.5 D.0.2")
        doc.add_paragraph("2. 下列哪种情况会使需求曲线右移？A.价格上升 B.收入增加 C.价格下降 D.替代品降价")
        doc.add_paragraph("二、计算题（15 分）")
        doc.add_paragraph("已知需求函数 Qd = 100 - 2P，供给函数 Qs = 20 + 2P，求均衡价格与均衡数量。")
        doc.add_paragraph("答案：令 Qd=Qs，得 100-2P=20+2P，P=20，Q=60。")
        doc.save(str(out_dir / "04_试题_往年真题.docx"))
        print(f"[sample-full] 已生成试题 {out_dir/'04_试题_往年真题.docx'}")
    except Exception as e:
        print(f"[sample-full] 跳过 exam docx（{e}）")

    print(f"[sample-full] 完成。可运行：py agent.py {out_dir} -o 笔记.md --course 微观经济学 --no-llm")


if __name__ == "__main__":
    main()
